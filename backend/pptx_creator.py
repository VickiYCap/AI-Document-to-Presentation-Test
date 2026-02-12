import json
import re
from pathlib import Path
from pptx import Presentation
import comtypes.client
import os

# Case-insensitive token pattern: {{ anything }}
TOKEN_REGEX = re.compile(r"{{(.*?)}}", re.IGNORECASE)

def find_all_tokens(presentation):
    """Return a sorted list of all token names found in the deck (without braces)."""
    tokens = set()
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for raw in TOKEN_REGEX.findall(run.text or ""):
                        tokens.add(raw.strip())
    return sorted(tokens)

def replace_tokens_in_presentation(presentation, token_map):
    """
    Replace tokens in the given Presentation using a flat dict token_map.
    Matching is case-insensitive.
    """
    # Normalize keys for case-insensitive lookup
    lower_map = {str(k).lower(): v for k, v in token_map.items()}
    replacements = 0

    for slide in presentation.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original = run.text or ""

                    def repl(m):
                        nonlocal replacements
                        key = (m.group(1) or "").strip().lower()
                        if key in lower_map:
                            replacements += 1
                            return str(lower_map[key])
                        # leave unknown tokens as-is
                        return m.group(0)

                    new_text = TOKEN_REGEX.sub(repl, original)
                    if new_text != original:
                        run.text = new_text

    return replacements

def load_flat_tokens(json_path_or_dict):
    """
    Accept either a dict or a path to a JSON file containing a flat token map.
    """
    if isinstance(json_path_or_dict, dict):
        return json_path_or_dict
    with open(json_path_or_dict, "r", encoding="utf-8") as f:
        return json.load(f)

#convert pptx to pdf for rendering  
def pptx_to_pdf(input_path):
    output_path = os.path.splitext(input_path)[0] + ".pdf"

    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    presentation = powerpoint.Presentations.Open(os.path.abspath(input_path))
    presentation.SaveAs(os.path.abspath(output_path), 32)
    presentation.Close()

    powerpoint.Quit()

    return output_path

def main():
    BASE = Path(__file__).resolve().parent
    powerpoints_dir = BASE.parent if (BASE / "powerpoints").exists() else BASE  # adjust if needed
    powerpoints_dir = BASE / "powerpoints"

    template_path = powerpoints_dir / "PPTX_Template.pptx"
    output_path   = powerpoints_dir / "Updated_Presentation.pptx"
    tokens_json   = powerpoints_dir / "eu_ai_policy_slides.json"  

    assert template_path.exists(), f"Template not found: {template_path}"
    assert tokens_json.exists(), f"Token map not found: {tokens_json}"

    # Load template
    pres = Presentation(str(template_path))

    # (Optional) Inspect tokens present in the deck
    present_tokens = find_all_tokens(pres)
    print("Tokens in template:", present_tokens)

    # Load flat token map and replace
    token_map = load_flat_tokens(tokens_json)
    replaced = replace_tokens_in_presentation(pres, token_map)

    # Report missing/unneeded
    present_lower = {t.lower() for t in present_tokens}
    provided_lower = {k.lower() for k in token_map.keys()}
    missing = sorted(present_lower - provided_lower)
    unused  = sorted(provided_lower - present_lower)

    # Save output (template remains untouched)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    pres.save(str(output_path))
    pptx_to_pdf(str(output_path))

    print(f"\nReplacements made: {replaced}")
    if missing:
        print("Tokens missing from JSON:", missing)
    if unused:
        print("Tokens in JSON but not in deck:", unused)
    print(f"Saved: {output_path}")

if __name__ == "__main__":
    main()