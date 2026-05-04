"""Analyzes the presentation template and then maps the JSON content to the correct shapes in the PPTX, 
including text, tables, and images. This is where the "magic" happens in terms of aligning the parsed 
content with the actual presentation structure. The mapping is heuristic-based, relying on text size, 
position, and shape type to guess which shape is the title, which are subtitles, body text, tables, and 
images. It then provides functions to fill in text (with bullet support), tables, and replace images
based on the JSON"""

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from typing import Dict, Any, List, Optional, Union
import json
import os
from copy import deepcopy
from pathlib import Path

# -----------------------------------------
# Check if a paragraph space is bulleted
# -----------------------------------------
def has_bullets(paragraph) -> bool:
    pPr = paragraph._p.pPr

    if pPr is not None:
        if pPr.find(qn('a:buChar')) is not None:
            return True
        if pPr.find(qn('a:buAutoNum')) is not None:
            return True
        if pPr.find(qn('a:buNone')) is not None:
            return False

    try:
        if paragraph.level and paragraph.level > 0:
            return True
    except AttributeError:
        pass

    return False


def build_body(shape):
    body = []
    if not getattr(shape, "has_text_frame", False):
        return body

    for paragraph in shape.text_frame.paragraphs:
        body.append({
            "bullets": has_bullets(paragraph)
        })

    return body

# -------------------------------------------------
# USE HEURISTIC TO DEFINE THE TITLE VS NORMAL TEXT
# -------------------------------------------------
def get_text_size_weight(shape, default_min_pt=12.0) -> float:
    if not getattr(shape, "has_text_frame", False):
        return float(default_min_pt)

    max_weight_pt = 0.0
    any_size_resolved = False

    for paragraph in shape.text_frame.paragraphs:
        para_size_len = paragraph.font.size 

        for run in paragraph.runs:
            size_len = run.font.size or para_size_len
            if size_len is None:
                continue

            any_size_resolved = True
            size_pt = float(size_len.pt)  
            factor = 2.0 if run.font.bold is True else 1.0
            max_weight_pt = max(max_weight_pt, size_pt * factor)

    if not any_size_resolved or max_weight_pt <= 0.0:
        max_weight_pt = float(default_min_pt)

    return max_weight_pt

# -------------------------------------------------
# cOUNT THE NUMBERS OF PARAGRAPHS PER SHAPE
# -------------------------------------------------
def count_paragraphs(shape) -> int:
    if not getattr(shape, "has_text_frame", False):
        return 0

    count = 0
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.text.strip():  
            count += 1

    return count

# -------------------------------------------------
# USE HEURISTIC TO GUESS WHAT THE TITLE IS IN A SLIDE
# -------------------------------------------------
def find_title(slide):
    candidates = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = (shape.text or "").strip()
        if not text:
            continue

        length = count_paragraphs(shape)
        if length > 1:
            continue
        top = shape.top
        if top <= 0:
            top = 0

        text_size_weight = get_text_size_weight(shape)
        score = (1 / (top + 1)) * 10000 * text_size_weight
        candidates.append({
            "shape": shape,
            "text": text,
            "score": score,
            "top": top,
            "text_size_weight": text_size_weight,
            "length" : length,
        })

    # print(candidates)

    if not candidates:
        return None
    best = max(candidates, key=lambda x: x["score"])
    # print("BEST: ", best)
    return best["shape"]

# -------------------------------------------------
# CHECK IF THERE IS A PICTURE PLACEHOLDER
# -------------------------------------------------
def is_picture_placeholder(shape) -> bool:
    if shape.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER:
        return False
    try:
        return shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE
    except Exception:
        return False

# ------------------------------------------------------------
# CHECK IF A PICTURE IS A BACKGROUND (SHOULDN'T BE REPLACED)
# ------------------------------------------------------------
def is_background_image(shape, slide_width, slide_height, threshold=0.45) -> bool:
    coverage = (shape.width * shape.height) / (slide_width * slide_height)
    return coverage >= threshold

# -------------------------------------------------------------------
# USING THE MAPPING FUNCTIONS CREATE THE JSON SCHEMA FOR ONE SLIDE
# -------------------------------------------------------------------
def slide_mapping(slide, slide_index=None, slide_width=None, slide_height=None):
    result = {"title": [], "subtitle": [], "body": [], "bullets": [], "table": [], "image": []}
    title_shape = find_title(slide)
    title_text = title_shape.text.strip() if title_shape and getattr(title_shape, "has_text_frame", False) else None

    for shape in slide.shapes:
        # --- image ---
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or is_picture_placeholder(shape):
            if slide_width and slide_height and is_background_image(shape, slide_width, slide_height):
                continue
            result["image"].append({
                "shape_id": shape.shape_id,
                "name": getattr(shape, "name", "") or "",
                "is_placeholder": shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER,
                "left_pct": shape.left / slide_width if slide_width else 0,
                "top_pct": shape.top / slide_height if slide_height else 0,
                "width_pct": shape.width / slide_width if slide_width else 1,
                "height_pct": shape.height / slide_height if slide_height else 1,
            })

        # --- tables ---
        if getattr(shape, "has_table", False):
            tb1 = shape.table
            n_rows = len(tb1.rows)
            n_cols = len(tb1.columns)

            result["table"].append({
                    "shape_id": shape.shape_id,
                    "name": getattr(shape, "name", "") or "",
                    "rows": n_rows,
                    "cols": n_cols,
                    "total_cells": n_rows * n_cols,
            })
            continue
            
        # --- text shapes ---
        if not getattr(shape, "has_text_frame", False):
            continue

        text = (shape.text or "").strip()
        if not text:
            continue

        if title_shape is not None and shape == title_shape:
            result["title"].append({
                "shape_id": shape.shape_id,
                "name": getattr(shape, "name", "") or ""
            })
            continue
        
        if count_paragraphs(shape) > 1 and shape.text != title_text:
            body_content = build_body(shape)

            if body_content:
                result["body"].append({
                    "shape_id": shape.shape_id,
                    "name": getattr(shape, "name", "") or "",
                    "content": body_content
                })

            continue


        result["subtitle"].append({
            "shape_id": shape.shape_id,
            "name": getattr(shape, "name", "") or ""
        })
    return result

# --------------------------------------------------------
# LOOP THROUGH PRESENTATION AND CREATE FULL JSON SCHEMA
# --------------------------------------------------------
def build_presentation_mapping(pptx_path):
    prs = Presentation(pptx_path)
    mapping = {}
    for idx, slide in enumerate(prs.slides, start=1):
        mapping[f"slide_{idx}"] = slide_mapping(slide, idx, prs.slide_width, prs.slide_height)
    return mapping

#----------------------------------------------------------- APPLYING TO PPTX ----------------------------------------------------------------------------------
# --------------------------------------------------------
#GET THE SHAPE ID FOR EACH SHAPE TO REPLACE TEXT/IMAGES
# --------------------------------------------------------
def find_shape_by_id(slide, shape_id) -> Optional[Any]:
    for shp in slide.shapes:
        if shp.shape_id == shape_id:
            return shp
    return None

# --------------------------------------------------------
# SET THE TEXT IN A TABLE CELL
# --------------------------------------------------------
def set_cell(cell, text: str):
    tf = cell.text_frame

    donor_rPr = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        donor_run = tf.paragraphs[0].runs[0]
        donor_rPr = deepcopy(donor_run._r.rPr) if donor_run._r.rPr is not None else None

    if len(tf.paragraphs) == 0:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]

    if len(p.runs) == 0:
        r = p.add_run()
    else:
        r = p.runs[0]

    r.text = str(text) if text is not None else ""

    if donor_rPr is not None:
        clone_run_rPr(r, donor_rPr)

    for extra in list(p.runs)[1:]:
        extra._r.getparent().remove(extra._r)

    for extra_p in list(tf.paragraphs)[1:]:
        tf._element.remove(extra_p._p)

# --------------------------------------------------------------
# cALLS SET_CELLS TO FILL IN THE ENTIRE TABLE WITH JSON CONTENT
# --------------------------------------------------------------
def fill_table(shape, rows_2d: List[Any], start_row: int = 0, pad_value: str = ""):
    if not hasattr(shape, "table"):
        return
    
    table = shape.table
    ppt_rows, ppt_cols = len(table.rows), len(table.columns)
    rows_2d = rows_2d or []

    def as_row_list(row):
        if row is None:
            vals = [pad_value]
        elif isinstance(row, (list, tuple)):
            vals = list(row)
        else:
            vals = [row]

        if len(vals) < ppt_cols:
            vals += [pad_value] * (ppt_cols - len(vals))
        elif len(vals) > ppt_cols:
            vals = vals[:ppt_cols]
        return [pad_value if v is None else str(v) for v in vals]
    
    data_i = 0
    for r in range(start_row, ppt_rows):
        row_vals = as_row_list(rows_2d[data_i]) if data_i < len(rows_2d) else [pad_value] * ppt_cols
        for c in range(ppt_cols):
            set_cell(table.cell(r, c), row_vals[c])
        data_i += 1

def clone_run_rPr(dst_run, donor_rPr):
    if donor_rPr is None:
        return
    dst_r = dst_run._r
    if dst_r.rPr is not None:
        dst_r.remove(dst_r.rPr)
    dst_r.insert(0, deepcopy(donor_rPr))

def set_body_text(shape, items):
    if not getattr(shape, "has_text_frame", False):
        return

    normalized_items = []
    for it in (items or []):
        if isinstance(it, str):
            normalized_items.append({"text": it, "bulleted": False})
        else:
            normalized_items.append({
                "text": it.get("text", ""),
                "bulleted": bool(it.get("bulleted", False)),
            })

    tf = shape.text_frame
    paragraphs = tf.paragraphs
    
    
    per_line_rPr = []
    last_style = None 

    for p in paragraphs:
        if p.runs:
            donor_run = p.runs[0]
            donor_rPr = deepcopy(donor_run._r.rPr) if donor_run._r.rPr is not None else None
            last_style = donor_rPr
        else:
            donor_rPr = deepcopy(last_style)  

        per_line_rPr.append(donor_rPr)
    
    base_p = paragraphs[0]
    base_pPr = deepcopy(base_p._p.get_or_add_pPr())

    for p in paragraphs[1:]:
        tf._element.remove(p._p)
    for r in list(base_p.runs):
        r._r.getparent().remove(r._r)

    for i, item in enumerate(normalized_items):
        text = item["text"]
        bulleted = item["bulleted"]

        para = base_p if i == 0 else tf.add_paragraph()

        cur_pPr = para._p.get_or_add_pPr()
        parent = para._p
        parent.remove(cur_pPr)
        parent.insert(0, deepcopy(base_pPr))

        try:
            para.level = 0
        except Exception:
            pass

        pPr = para._p.get_or_add_pPr()
        for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buBlip"):
            el = pPr.find(qn(tag))
            if el is not None:
                pPr.remove(el)

        if bulleted:
            buChar = OxmlElement("a:buChar")
            buChar.set("char", "•")
            pPr.append(buChar)
        else:
            pPr.append(OxmlElement("a:buNone"))

        r = para.add_run()
        clone_run_rPr(r, per_line_rPr[i])
        r.text = item["text"]


# --------------------------------------------------------
# SIMPLY ADDS THE PLAIN TEXT INTO THE NEW PPTX
# --------------------------------------------------------
def set_plain_text(shape, text: str):
    if not getattr(shape, "has_text_frame", False):
        return

    tf = shape.text_frame
    if len(tf.paragraphs) == 0:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]

    if len(p.runs) == 0:
        r = p.add_run()
    else:
        r = p.runs[0]
    r.text = text

    for extra in list(p.runs)[1:]:
        extra._r.getparent().remove(extra._r)

    for extra_p in list(tf.paragraphs)[1:]:
        tf._element.remove(extra_p._p)

# --------------------------------------------------------
# HANDLES IMAGE REPLACEMENT IN THE PPTX, WORKS FOR BOTH
#  PICTURE PLACEHOLDERS AND REGULAR IMAGES
# --------------------------------------------------------
def replace_image(shape, path: str):
    if not os.path.isfile(path):
        raise FileNotFoundError(f"Image not found: {path}")

    sp_tree = shape._element.getparent()
    old_idx = list(sp_tree).index(shape._element)

    if is_picture_placeholder(shape):
        try:
            return shape.insert_picture(path)
        except Exception:
            slide = shape.part.slide
            new_pic = slide.shapes.add_picture(path, shape.left, shape.top, width=shape.width, height=shape.height)
            try:
                sp_tree.remove(shape._element)
            except Exception:
                pass
            sp_tree.remove(new_pic._element)
            sp_tree.insert(old_idx, new_pic._element)
            return new_pic

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        slide = shape.part.slide
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        new_pic = slide.shapes.add_picture(path, left, top, width=width, height=height)
        sp_tree.remove(shape._element)
        sp_tree.remove(new_pic._element)
        sp_tree.insert(old_idx, new_pic._element)
        return new_pic

    raise TypeError("Target shape is neither a picture placeholder nor a picture shape.")

# --------------------------------------------------------
# APPLIES THE IMAGE REPLACEMENTS ONTO THE NEW SLIDES
# --------------------------------------------------------
def apply_images(slide, slide_map: Dict[str, Any], images_spec: Union[List[str], Dict[str, Any]]):
    img_targets = slide_map.get("image", [])
    if not img_targets:
        return

    if isinstance(images_spec, list):
        for i, path in enumerate(images_spec):
            if i >= len(img_targets):
                break
            if not path or not os.path.isfile(path):
                continue
            ref = img_targets[i]
            shp = find_shape_by_id(slide, ref["shape_id"])
            if shp is not None:
                replace_image(shp, path)
        return
    if isinstance(images_spec, dict):
        mode = images_spec.get("by")
        items = images_spec.get("items", [])
        if mode == "id":
            for item in items:
                sid = item.get("shape_id")
                path = item.get("path")
                if sid is None or not path:
                    continue
                shp = find_shape_by_id(slide, sid)
                if shp is not None:
                    replace_image(shp, path)
        elif mode == "name":
            name_map = {}
            for ref in img_targets:
                name_map[ref.get("name", "")] = find_shape_by_id(slide, ref["shape_id"])
            for item in items:
                name = item.get("name", "")
                path = item.get("path")
                shp = name_map.get(name)
                if shp is not None and path:
                    replace_image(shp, path)


# ---------------------------------------------------------------------
# APPLIES TEXT AND IMAGES FROM THE JSON CREATOR TO THE NEW PPTX SLIDES
# ---------------------------------------------------------------------
def apply_json_to_slide(slide, slide_map: Dict[str, Any], json_slide: Dict[str, Any]):
    # 0) Images (do early so content doesn't overlap text work)
    if "images" in json_slide:
        apply_images(slide, slide_map, json_slide["images"])

    # 1) Title
    if "title" in json_slide and slide_map.get("title"):
        title_shape_ref = slide_map["title"][0]   # first title only
        shp = find_shape_by_id(slide, title_shape_ref["shape_id"])
        if shp is not None:
            set_plain_text(shp, json_slide["title"])

    # 2) Regular text (map in order to "regular" shapes)
    if "subtitle" in json_slide and slide_map.get("subtitle"):
        values = json_slide["subtitle"]
        regs = slide_map["subtitle"]
        for idx, value in enumerate(values):
            if idx >= len(regs):
                break
            shp = find_shape_by_id(slide, regs[idx]["shape_id"])
            if shp is not None:
                set_plain_text(shp, value)

     # 2) body text
    if "body" in json_slide and slide_map.get("body"):
        values = json_slide["body"]
        regs = slide_map["body"]

        for idx, value_list in enumerate(values):
            if idx >= len(regs):
                break

            shp = find_shape_by_id(slide, regs[idx]["shape_id"])
            if shp is None:
                continue
            set_body_text(shp, value_list)

    # 4) Tables
    if "table" in json_slide and slide_map.get("table"):
        table_spec = json_slide["table"]
        table_shape_ref = slide_map["table"][0] if isinstance(slide_map["table"], list) else slide_map["table"]
        shp = find_shape_by_id(slide, table_shape_ref["shape_id"])
        if shp is not None:

            def extract_table_rows_2d(obj):
                if obj is None:
                    return []

                if isinstance(obj, list) and obj and isinstance(obj[0], dict) and "rows" in obj[0]:
                    return obj[0].get("rows", []) or []

                if isinstance(obj, dict) and "rows" in obj:
                    return obj.get("rows", []) or []

                if isinstance(obj, list) and obj and all(isinstance(x, list) for x in obj):
                    return obj

                if isinstance(obj, list):
                    return [[x] for x in obj]

                return []

            data_rows = extract_table_rows_2d(table_spec)
            fill_table(shp, data_rows)

# -------------------------------------------------
# MAPS THE JSON FILLER TO THE ACTUAL PPTX
# -------------------------------------------------
def apply_json_to_pptx(pptx_in: str, json_in: str, pptx_out: str, mapping: Optional[Dict[str, Any]] = None):
    prs = Presentation(pptx_in)
    if mapping is None:
        mapping = build_presentation_mapping(pptx_in)

    with open(json_in, "r", encoding="utf-8") as f:
        data = json.load(f)

    slides_json = data.get("slides", [])
    total_slides = min(len(slides_json), len(prs.slides))

    for idx in range(total_slides):
        slide = prs.slides[idx]
        slide_key = f"slide_{idx+1}"
        slide_map = mapping.get(slide_key, None)
        if not slide_map:
            continue
        apply_json_to_slide(slide, slide_map, slides_json[idx])

    prs.save(pptx_out)

# -------------------------------------------------
# Testing function to preview JSON (unused)
# -------------------------------------------------
def preview_alignment(pptx_path: str, json_path: str):
    """
    Print a quick alignment summary without modifying files:
    - Number of slides
    - Which slides have title/regular/bullets/body/table/img targets vs. JSON content
    """
    mapping = build_presentation_mapping(pptx_path)
    prs = Presentation(pptx_path)
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    slides_json = data.get("slides", [])

    print(f"PPTX slides: {len(prs.slides)} | JSON slides: {len(slides_json)}")
    for idx, js in enumerate(slides_json, start=1):
        sm = mapping.get(f"slide_{idx}", {})
        print(f"\nSlide {idx}:")
        print(f"  JSON keys: {list(js.keys())}")
        print(
            f"  Mapped targets: title={len(sm.get('title', []))}, "
            f"subtitle={len(sm.get('subtitle', []))}, bullets={len(sm.get('bullets', []))}, "
            f"body={len(sm.get('body', []))}, table={len(sm.get('table', []))}, "
            f"img={len(sm.get('image', []))}"
        )

if __name__ == "__main__":
    # prs = Presentation("pptx_templates/Cap_Template.pptx")
    # slide = prs.slides[5]
    # find_title(slide)
    
    template = "pptx_templates/Cap_Template.pptx"
    template_path = Path(__file__).resolve().parent / template
    pptx_schema = build_presentation_mapping(template_path)
    print (pptx_schema)

    # preview_alignment(
    #     "pptx_templates/Heuristic_Template.pptx",
    #     "pptx_templates/heuristic.json"
    # )

    base = Path(__file__).resolve().parent
    apply_json_to_pptx(
        pptx_in=str(base / "pptx_templates/Cap_Template.pptx"),
        json_in=str(base / "pptx_templates/Cap_Filled.json"),
        pptx_out=str(base / "pptx_templates/Cap_Filled.pptx")
    )
    print("Done. Saved: pptx_templates/Cap_Filled.pptx")