from pptx import Presentation
from pptx.oxml.ns import qn
import json

# Titles → big font, near top, usually 1 line, sometimes bold
# Bulleted body text → smaller font, has bullets, multiple lines
# Regular body text → smaller font, no bullets, usually short

def has_bullets(shape):
    if not shape.has_text_frame:
        return False

    for paragraph in shape.text_frame.paragraphs:
        pPr = paragraph._p.pPr  
        if pPr is None:
            continue

        # Check for bullet characters (•, –, etc.)
        bu_char = pPr.find(qn('a:buChar'))
        # Check for numbered lists (1., a., i., etc.)
        bu_auto = pPr.find(qn('a:buAutoNum'))

        if bu_char is not None or bu_auto is not None:
            return True

    return False

def get_text_size_weight(shape):
    bold = 0
    if not shape.has_text_frame:
        return 0
    
    # print("Shape:", shape.text.strip())
    # for p in shape.text_frame.paragraphs:
    #     for r in p.runs:
    #         print("   Run:", repr(r.text), "size:", r.font.size)

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None and run.font.bold:
                bold = 1
    return run.font.size.pt * bold  

def count_paragraphs(shape):
    if not getattr(shape, "has_text_frame", False):
        return 0
    return int(len(shape.text_frame.paragraphs))


def find_title(slide):
    candidates = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        text = shape.text.strip()
        if len(text) == 0:
            continue

        # Extract geometry
        length = count_paragraphs(shape)
        top = shape.top
        text_size_weight = get_text_size_weight(shape)


        # Heuristic score:
        # - higher if closer to the top
        # - higher if wider/taller

        # top is a smaller number if its higher in the slide
        score = (1 / (top + 1)) * 10000 + text_size_weight + (1/length)

        candidates.append({
            "shape": shape,
            "text": text,
            "score": score,
            "text_size_weight": text_size_weight,
        })
    # for c in candidates:
    #     print(f"Candidate: '{c['text']}' Score: {c['score']} weight: {c['text_size_weight']}")

    if not candidates:
        return None

    # Pick the highest-scoring candidate
    best = max(candidates, key=lambda x: x["score"])
    #print("title: ", best["text"])
    return best["shape"]

def replace_title(slide, new_title):
    title_shape = find_title(slide)
    if title_shape:
        title_shape.text = new_title

def find_shapes(slide):
    title = find_title(slide)
    print("TITLE:", title.text)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        text = shape.text.strip()
        if len(text) == 0:
            continue

        if has_bullets(shape):
            print("BULLET BOX:", shape.text)
        
        if not has_bullets(shape) and shape.text != title.text:
            if count_paragraphs(shape) > 1:
                print("BODY TEXT:", shape.text)
            else:
                print("REGULAR TEXT", shape.text)

prs = Presentation("powerpoints/Heuristic_Template.pptx")

slide = prs.slides[0]
#find_shapes(slide)
find_title(slide)

for i in range(len(prs.slides)):
    print(f"Slide {i+1}:")
    find_shapes(prs.slides[i])
    print("\n")

